"""Generate a 3-strategy presentation deck for the Bachelor thesis defense."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
C_ACCENT_A  = RGBColor(0x00, 0xB4, 0xD8)   # cyan  – Strategy A
C_ACCENT_B  = RGBColor(0xFF, 0x6B, 0x6B)   # coral – Strategy B
C_ACCENT_C  = RGBColor(0x06, 0xD6, 0xA0)   # green – Strategy C
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
C_DIM       = RGBColor(0x88, 0x99, 0xAA)
C_DARK_CARD = RGBColor(0x1A, 0x2E, 0x42)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    color = color or C_WHITE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)


def accent_bar(slide, color, height=Inches(0.06)):
    add_rect(slide, 0, 0, SLIDE_W, height, color)


def add_slide_number(slide, n, color):
    add_text_box(slide, str(n), Inches(12.7), Inches(7.1), Inches(0.5), Inches(0.3),
                 font_size=10, color=color, align=PP_ALIGN.RIGHT)


def add_section_label(slide, label, color):
    add_text_box(slide, label, Inches(0.35), Inches(0.12), Inches(6), Inches(0.3),
                 font_size=10, bold=True, color=color)


def bullet_frame(slide, items, x, y, w, h, font_size=15, color=None, bullet="•  "):
    color = color or C_WHITE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ── Slide factories ────────────────────────────────────────────────────────

def make_title_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide)

    # decorative accent bars per strategy
    add_rect(slide, 0, Inches(2.8), Inches(4.44), Inches(0.05), C_ACCENT_A)
    add_rect(slide, Inches(4.44), Inches(2.8), Inches(4.44), Inches(0.05), C_ACCENT_B)
    add_rect(slide, Inches(8.88), Inches(2.8), Inches(4.44), Inches(0.05), C_ACCENT_C)

    add_text_box(slide, "ML Experiment Optimization with Intelligent Caching",
                 Inches(1), Inches(1.0), Inches(11.33), Inches(1.2),
                 font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Bachelor Thesis Presentation  |  Lennart Gorzel  |  IMC FH Krems",
                 Inches(1), Inches(2.2), Inches(11.33), Inches(0.5),
                 font_size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Three presentation strategies — pick one, mix as needed",
                 Inches(1), Inches(3.1), Inches(11.33), Inches(0.4),
                 font_size=12, color=C_DIM, align=PP_ALIGN.CENTER, italic=True)

    labels = [
        ("Strategy A", "The Researcher's Day", C_ACCENT_A),
        ("Strategy B", "The Trust Problem", C_ACCENT_B),
        ("Strategy C", "Numbers-First", C_ACCENT_C),
    ]
    for i, (tag, subtitle, col) in enumerate(labels):
        bx = Inches(0.5 + i * 4.27)
        add_rect(slide, bx, Inches(3.8), Inches(3.9), Inches(1.5), C_DARK_CARD)
        add_text_box(slide, tag, bx + Inches(0.15), Inches(3.9), Inches(3.6), Inches(0.5),
                     font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text_box(slide, subtitle, bx + Inches(0.15), Inches(4.45), Inches(3.6), Inches(0.5),
                     font_size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Slides 2–7: Strategy A  |  Slides 9–14: Strategy B  |  Slides 16–21: Strategy C",
                 Inches(1), Inches(6.7), Inches(11.33), Inches(0.4),
                 font_size=10, color=C_DIM, align=PP_ALIGN.CENTER)


def make_divider(prs, label, subtitle, color, slide_num):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, color)
    add_rect(slide, Inches(0.15), Inches(2.9), SLIDE_W - Inches(0.15), Inches(0.04), color)
    add_text_box(slide, label, Inches(1), Inches(1.4), Inches(11), Inches(1.2),
                 font_size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, subtitle, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
                 font_size=18, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    timing_text = "Timing per section:  Motivation 1:30  •  SOTA 0:45  •  Model 1:00  •  Implementation 1:15  •  Results 1:30  •  Conclusion 1:00" \
        if label.startswith("Strategy A") else \
        "Timing per section:  Motivation 1:00  •  SOTA 0:45  •  Model 1:15  •  Implementation 1:15  •  Results 1:45  •  Conclusion 1:00" \
        if label.startswith("Strategy B") else \
        "Timing per section:  Motivation 1:15  •  SOTA 0:45  •  Model 1:00  •  Implementation 1:15  •  Results 1:30  •  Conclusion 1:15"
    add_text_box(slide, timing_text, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
                 font_size=11, color=C_DIM, align=PP_ALIGN.CENTER)

    add_text_box(slide, "7 min  |  CS committee  |  English",
                 Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.35),
                 font_size=12, color=C_DIM, align=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num, color)
    return slide


def make_content_slide(prs, section_tag, title, body_items, notes_lines,
                       color, slide_num, hook=None, time_budget=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    accent_bar(slide, color)

    # Section tag top-left
    add_section_label(slide, section_tag, color)

    # Time budget chip top-right
    if time_budget:
        add_text_box(slide, time_budget, Inches(11.5), Inches(0.12), Inches(1.6), Inches(0.3),
                     font_size=10, color=color, align=PP_ALIGN.RIGHT, bold=True)

    # Title
    add_text_box(slide, title, Inches(0.35), Inches(0.45), Inches(12.6), Inches(0.75),
                 font_size=26, bold=True, color=C_WHITE)

    # Divider line under title
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(12.6), Inches(0.02), color)

    # Hook callout (if any)
    content_top = Inches(1.4)
    if hook:
        hook_box = add_rect(slide, Inches(0.35), Inches(1.4), Inches(12.6), Inches(0.75), C_DARK_CARD)
        add_text_box(slide, f'"{hook}"', Inches(0.55), Inches(1.44), Inches(12.2), Inches(0.68),
                     font_size=14, italic=True, color=color, align=PP_ALIGN.CENTER)
        content_top = Inches(2.25)

    # Body bullets
    bullet_frame(slide, body_items, Inches(0.35), content_top,
                 Inches(8.4), Inches(7.5) - content_top - Inches(0.5),
                 font_size=15, color=C_WHITE)

    # Speaker notes panel
    note_y = Inches(5.6)
    add_rect(slide, Inches(8.9), Inches(1.4), Inches(4.1), Inches(5.8), C_DARK_CARD)
    add_text_box(slide, "SPEAKER NOTES", Inches(9.0), Inches(1.5), Inches(3.9), Inches(0.3),
                 font_size=9, bold=True, color=color)
    add_rect(slide, Inches(9.0), Inches(1.85), Inches(3.8), Inches(0.01), color)
    note_text = "\n".join(f"• {n}" for n in notes_lines)
    txBox = slide.shapes.add_textbox(Inches(9.0), Inches(1.95), Inches(3.8), Inches(5.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for n in notes_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "• " + n
        run.font.size = Pt(11)
        run.font.color.rgb = C_LIGHT

    add_slide_number(slide, slide_num, color)
    return slide


# ── Content data ───────────────────────────────────────────────────────────

def slides_strategy_a(prs, start_n):
    col = C_ACCENT_A
    tag_prefix = "Strategy A — The Researcher's Day"

    slides = [
        dict(
            section_tag=f"{tag_prefix}  |  Motivation",
            title="Motivation — Nine Hours. Again.",
            hook="I run an experiment. 9 hours. I change one number. Another 9 hours.",
            body_items=[
                "Sleep EEG classification — 128 subjects, ~107,000 epochs",
                "149 features/epoch  ×  128 LOSO folds  ×  18 configurations = 2,304 training runs",
                "One full experiment run: 9–12 hours",
                "Tweak one hyperparameter → start over → same 9–12 hours",
                "Over 80 % of compute is identical recomputation",
                "→ The bottleneck is not the model — it's the pipeline",
            ],
            notes=[
                "Open with the hook line verbatim — pause after it.",
                "Let the 2,304 number land. Write it on screen, say it once.",
                "Don't explain LOSO yet — keep it as 'standard cross-validation'.",
                "End the section with: 'This isn't a model problem. It's a pipeline problem.'",
            ],
            time="1:30",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  State of the Art",
            title="State of the Art — The Parts Exist",
            body_items=[
                "Cryptographic hashing (SHA-256) — deterministic, collision-resistant ✓",
                "Caching frameworks — joblib, functools.lru_cache, MLflow artifact stores ✓",
                "Model serialization — pickle, joblib.dump, ONNX ✓",
                "────────────────────────────────────────────",
                "→ Each handles one piece.",
                "→ None handles the boundary between them under LOSO cross-validation.",
                "→ Change one parameter and you cannot know which cached results are still valid.",
            ],
            notes=[
                "Three columns, one line each — this slide should feel fast.",
                "DO NOT mention your combination here. Just show the gap.",
                "The last three lines are the gap. Read them slowly.",
                "Transition: 'The parts exist. Assembling them safely does not.'",
            ],
            time="0:45",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Model",
            title="Model — A Computation Has an Identity",
            body_items=[
                "Core idea: every computation has an identity derived from all its inputs",
                "Identity = SHA-256( all inputs ) → cache key",
                "If nothing changed → key matches → load from cache",
                "If anything changed → key differs → recompute",
                "Downstream stages inherit upstream fingerprints",
                "→ Partial invalidation: change feature selection, keep preprocessing cache",
            ],
            notes=[
                "Show one small diagram: [inputs] → [SHA-256] → [cache key] → [output]",
                "The word 'identity' is the anchor concept — return to it.",
                "Avoid the word 'hash function' if the audience might not know it — use 'fingerprint'.",
                "Transition: 'That's the concept. Now the wiring.'",
            ],
            time="1:00",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Implementation",
            title="Implementation — 4-Stage Pipeline",
            body_items=[
                "Stage 1 — Preprocessing:  fingerprint = raw EEG + filter config",
                "Stage 2 — Feature Extraction:  fingerprint = preprocessed signals + extraction config",
                "Stage 3 — Feature Selection:  fingerprint = features + corr threshold + top-K",
                "Stage 4 — Model Training:  fingerprint = selected features + model params + subject ID",
                "────────────────────────────────────────────",
                "Subject ID in Stage 4 fingerprint → LOSO fold identity is part of the cache key",
                "Change config at Stage 3 → only Stage 3 + 4 recompute → Stages 1 & 2 reused",
            ],
            notes=[
                "Show one pipeline diagram with arrows — no code, no file tree.",
                "The subject-ID point is the key novelty. Say it once, clearly.",
                "Demonstrate partial invalidation with a concrete example.",
                "Keep this slide visual. The diagram does the talking.",
            ],
            time="1:15",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Results",
            title="Results — Before and After",
            body_items=[
                "Cold run (no cache):       ~9 hours",
                "Warm run (full cache):    ~1 hour          → 4.5× overall speedup",
                "Feature extraction alone: 53 min → 0.2 min  → 224× speedup",
                "Cache hit rate:            100 %  (24/24 models, 0 errors)",
                "Data leakage:              0 confirmed cases",
                "────────────────────────────────────────────",
                "Iteration cycle: 9h → 1h  means ~9 experiments per day instead of 1",
            ],
            notes=[
                "Use a bar chart — two bars (cold / warm) — not a table.",
                "The 224× is the showstopper number. Pause after it.",
                "Lead with speed, then confirm correctness (100% hit, 0 leakage).",
                "The last line reframes the impact: researcher velocity, not just CPU time.",
            ],
            time="1:30",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Conclusion",
            title="Conclusion — The Bottleneck Is Gone",
            body_items=[
                "Fingerprint everything. Cache everything. Recompute only what changed.",
                "9-hour iteration cycle → 1-hour iteration cycle",
                "Safe: subject ID in fingerprint prevents leakage by construction",
                "Generalizes to any LOSO / k-fold pipeline with expensive features",
                "────────────────────────────────────────────",
                "\"Iteration cost was the bottleneck of this research. It isn't anymore.\"",
            ],
            notes=[
                "Closing line is the thesis in one sentence — deliver it without looking at the slide.",
                "Mention generalization briefly: 'any expensive ML pipeline under k-fold.'",
                "Keep it short — 60 seconds. The results already made the argument.",
                "End with a pause, then open for questions.",
            ],
            time="1:00",
        ),
    ]

    for i, s in enumerate(slides):
        make_content_slide(prs,
                           section_tag=s["section_tag"],
                           title=s["title"],
                           body_items=s["body_items"],
                           notes_lines=s["notes"],
                           color=col,
                           slide_num=start_n + i,
                           hook=s.get("hook"),
                           time_budget=s["time"])


def slides_strategy_b(prs, start_n):
    col = C_ACCENT_B
    tag_prefix = "Strategy B — The Trust Problem"

    slides = [
        dict(
            section_tag=f"{tag_prefix}  |  Motivation",
            title="Motivation — Speed Is Easy. Safety Isn't.",
            hook="Cross-validation has a rule: don't touch the held-out subject. Caching has a habit: it remembers everything.",
            body_items=[
                "Sleep EEG: 128 subjects, 2,304 training runs, 9–12 h per run",
                "80 %+ of compute is identical across runs — obvious to cache",
                "But: LOSO cross-validation requires strict data isolation per fold",
                "Wrong cache → wrong held-out subject in training → silent data leakage",
                "→ The real problem is not speed — it is caching without risking correctness",
            ],
            notes=[
                "Read the hook line verbatim. Let the contradiction land.",
                "The audience knows why cross-validation matters — use that.",
                "Frame the complication as a trust problem, not just a perf problem.",
                "Transition: 'So the question is: can we cache safely?'",
            ],
            time="1:00",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  State of the Art",
            title="State of the Art — Building Blocks, No Glue",
            body_items=[
                "SHA-256 hashing — deterministic fingerprinting of arbitrary data ✓",
                "Caching libraries — joblib, MLflow, lru_cache ✓",
                "Model serialization — pickle, joblib.dump, ONNX ✓",
                "────────────────────────────────────────────",
                "→ None encodes LOSO fold identity into the cache key.",
                "→ None provides stage-wise partial invalidation.",
                "→ The gap: no tool ties cache validity to cross-validation fold boundaries.",
            ],
            notes=[
                "Keep this fast — one sentence per tool.",
                "The gap is in the last three lines. Slow down there.",
                "Do NOT reveal your solution here. Point at the empty space, walk away.",
                "Transition: 'The parts exist. The boundary between them under cross-validation does not.'",
            ],
            time="0:45",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Model",
            title="Model — Identity Includes the Fold",
            body_items=[
                "A computation's identity = SHA-256 of ALL its inputs",
                "For Stage 4 (model training), inputs include: features + params + held-out subject ID",
                "→ Each LOSO fold gets a unique cache key by construction",
                "→ Cached model for fold 42 cannot be loaded for fold 17",
                "→ Data leakage is prevented at the identity level, not by runtime checks",
                "Downstream stages inherit upstream fingerprints → hierarchical invalidation",
            ],
            notes=[
                "This is the insight slide. Give it 15 extra seconds if needed.",
                "Emphasize 'by construction' — it's mathematically guaranteed, not guarded.",
                "Diagram: [inputs + subject_id] → [SHA-256] → [cache key]",
                "Transition: 'That's the concept. Now the wiring.'",
            ],
            time="1:15",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Implementation",
            title="Implementation — 4 Stages, Each Fingerprinted",
            body_items=[
                "Stage 1 — Preprocessing:   key = raw data + filter settings",
                "Stage 2 — Features:         key = preprocessed + extraction config",
                "Stage 3 — Selection:         key = features + threshold + top-K",
                "Stage 4 — Model:             key = selected features + params + subject ID  ← fold guard",
                "────────────────────────────────────────────",
                "Partial invalidation: change threshold → only Stages 3 & 4 recompute",
                "Cache validated: 100 % hit rate, 0 serialization errors, 0 leakage events",
            ],
            notes=[
                "Show the pipeline diagram. Highlight Stage 4's subject ID addition.",
                "The partial invalidation line shows the system is hierarchical, not monolithic.",
                "Don't show code. The diagram is cleaner and faster.",
                "Transition: 'Does it hold up? Let the numbers answer.'",
            ],
            time="1:15",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Results",
            title="Results — Fast AND Correct",
            body_items=[
                "Speed:     9 h (cold) → 1 h (warm)  — 4.5× overall speedup",
                "Features:  53 min → 0.2 min           — 224× speedup",
                "────────────────────────────────────────────",
                "Correctness:  100 % cache hit rate (24/24 models loaded correctly)",
                "Leakage:      0 confirmed cases across all 128 LOSO folds",
                "────────────────────────────────────────────",
                "→ The two claims hold simultaneously: faster AND safe",
            ],
            notes=[
                "Present speed first, then correctness — that's the pair this strategy promises.",
                "The '0 leakage' line is the payoff for the trust framing. Make it land.",
                "Bar chart: cold vs warm. Then a second visual: hit rate + leakage table.",
                "Don't rush here — this is the proof of both claims.",
            ],
            time="1:45",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Conclusion",
            title="Conclusion — Caching That Survives Cross-Validation",
            body_items=[
                "Caching is well-understood. Caching safely under LOSO was not.",
                "Solution: encode fold identity into the cache key — prevention by construction",
                "Result: 80 %+ speedup with zero correctness compromise",
                "Generalizes to any k-fold pipeline with expensive intermediate computations",
                "────────────────────────────────────────────",
                "\"Caching, safely, under cross-validation.\"",
            ],
            notes=[
                "Open with the contrast: known vs. unknown — that's the contribution.",
                "The closing one-liner is the thesis. Say it slowly, then stop.",
                "Mention generalization to invite follow-up questions if time allows.",
                "Don't add slides. Don't summarize the slides. Just land the message.",
            ],
            time="1:00",
        ),
    ]

    for i, s in enumerate(slides):
        make_content_slide(prs,
                           section_tag=s["section_tag"],
                           title=s["title"],
                           body_items=s["body_items"],
                           notes_lines=s["notes"],
                           color=col,
                           slide_num=start_n + i,
                           hook=s.get("hook"),
                           time_budget=s["time"])


def slides_strategy_c(prs, start_n):
    col = C_ACCENT_C
    tag_prefix = "Strategy C — Numbers-First"

    slides = [
        dict(
            section_tag=f"{tag_prefix}  |  Motivation",
            title="Motivation — 9 Hours. Then 1 Hour.",
            hook="9 hours. Then 1 hour. Same experiment, same data, same correctness. One design choice.",
            body_items=[
                "Sleep EEG classification: 128 subjects, LOSO cross-validation",
                "128 folds × 18 configurations = 2,304 training runs",
                "One run: 9–12 hours end-to-end",
                "80 %+ of that time: identical recomputation across runs",
                "→ How do you eliminate redundant compute without risking wrong results?",
            ],
            notes=[
                "Open with the hook number before any context. Let curiosity precede explanation.",
                "Then backfill: here's what the 9 hours buys you, here's why there are so many.",
                "End with the question — plant it early so every following slide answers it.",
                "This is a result-first strategy: the audience now knows the payoff.",
            ],
            time="1:15",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  State of the Art",
            title="State of the Art — Close, But Not Quite",
            body_items=[
                "Hashing (SHA-256)           — fingerprints arbitrary data  ✓",
                "Caching frameworks           — joblib, MLflow, lru_cache  ✓",
                "Model serialization          — pickle, ONNX, joblib.dump  ✓",
                "────────────────────────────────────────────",
                "Missing: a unified scheme that ties all three together under LOSO boundaries",
                "→ None prevents stale-fold cache reads. None cascades invalidation across stages.",
            ],
            notes=[
                "Fast slide — 45 seconds, three tools, one gap.",
                "The gap line justifies why there was a thesis to write.",
                "Don't linger here. The audience already knows the answer is coming.",
            ],
            time="0:45",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Model",
            title="Model — Fingerprint Every Computation",
            body_items=[
                "Every computation has an identity: SHA-256( all inputs )",
                "Same inputs → same key → load from cache",
                "Different inputs → different key → recompute",
                "Stages are chained: downstream fingerprint includes upstream fingerprint",
                "→ Change one stage → only that stage and all downstream stages invalidate",
            ],
            notes=[
                "One diagram: inputs → SHA-256 → key → cache.",
                "Emphasize the chaining — it's what makes partial invalidation possible.",
                "Keep conceptual, no implementation details yet.",
                "Transition: 'That's the idea. Here's the wiring.'",
            ],
            time="1:00",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Implementation",
            title="Implementation — 4 Stages, Hierarchically Cached",
            body_items=[
                "Stage 1 — Preprocessing:   fingerprint = raw EEG + filter settings",
                "Stage 2 — Features:         fingerprint = Stage 1 output + extraction config",
                "Stage 3 — Selection:         fingerprint = Stage 2 output + threshold + top-K",
                "Stage 4 — Model:             fingerprint = Stage 3 output + params + subject ID",
                "────────────────────────────────────────────",
                "Subject ID in Stage 4 → each LOSO fold is a distinct cache entry → no leakage",
                "Change threshold → only Stages 3 + 4 recompute → ~8× further speedup",
            ],
            notes=[
                "Show the 4-stage diagram with arrows.",
                "Point at Stage 4 and say 'this is where the safety guarantee lives.'",
                "One concrete example of partial invalidation — corr threshold 0.75 → 0.90.",
                "Transition: 'Now — does it actually deliver on the promise?'",
            ],
            time="1:15",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Results",
            title="Results — The Numbers Hold Up",
            body_items=[
                "Overall:    9 h (cold) → 1 h (warm)           = 4.5× speedup",
                "Features:   53 min → 0.2 min                    = 224× speedup",
                "Hit rate:   100 %  (24/24 models, 0 errors)",
                "Leakage:    0 confirmed cases (128 LOSO folds)",
                "────────────────────────────────────────────",
                "Expected at full scale: 8–10× overall speedup",
                "→ The promise from slide 1 is delivered — same correctness, 9× faster",
            ],
            notes=[
                "The audience has been waiting for these since slide 1. Don't underdeliver.",
                "Bar chart: cold vs warm. It's visual and fast.",
                "Explicitly echo the opening hook: 'Remember slide 1 — 9 hours then 1 hour.'",
                "The last bullet closes the loop on the opening question.",
            ],
            time="1:30",
        ),
        dict(
            section_tag=f"{tag_prefix}  |  Conclusion",
            title="Conclusion — One Design Choice",
            body_items=[
                "Fingerprint every computation with all its inputs — including the fold ID",
                "4 cascading stages → partial invalidation → cache only what changed",
                "Result: 4.5–10× speedup with no correctness compromise",
                "Generalizes: any k-fold pipeline with expensive feature extraction",
                "────────────────────────────────────────────",
                "\"9 hours. Then 1 hour. Same experiment. One design choice.\"",
            ],
            notes=[
                "Echo the opening line as the closing line — that's the story arc complete.",
                "Mention generalization: EEG-specific in implementation, general in principle.",
                "Stop after the closing line. Don't add a summary of the summary.",
                "Pause. Wait for questions.",
            ],
            time="1:15",
        ),
    ]

    for i, s in enumerate(slides):
        make_content_slide(prs,
                           section_tag=s["section_tag"],
                           title=s["title"],
                           body_items=s["body_items"],
                           notes_lines=s["notes"],
                           color=col,
                           slide_num=start_n + i,
                           hook=s.get("hook"),
                           time_budget=s["time"])


# ── Build ──────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: master title
    make_title_slide(prs)

    # Slide 2: Strategy A divider
    make_divider(prs,
                 "Strategy A — The Researcher's Day",
                 "Narrative: first-person pain → unified idea → relief",
                 C_ACCENT_A, 2)
    # Slides 3–8: Strategy A
    slides_strategy_a(prs, start_n=3)

    # Slide 9: Strategy B divider
    make_divider(prs,
                 "Strategy B — The Trust Problem",
                 "Narrative: correctness framing → gap → proof of speed AND safety",
                 C_ACCENT_B, 9)
    # Slides 10–15: Strategy B
    slides_strategy_b(prs, start_n=10)

    # Slide 16: Strategy C divider
    make_divider(prs,
                 "Strategy C — Numbers-First",
                 "Narrative: open with the big number, backfill the why, close the loop",
                 C_ACCENT_C, 16)
    # Slides 17–22: Strategy C
    slides_strategy_c(prs, start_n=17)

    out = "/home/user/BachlorThesis/presentation_strategies.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
