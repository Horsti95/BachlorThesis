"""V2: 5x5-rule presentation deck. Max 5 bullets per slide, max ~5 words each.

Speaker notes go into PowerPoint's notes pane (not on the slide face).
Embeds pipeline figure (Implementation) and speedup bar chart (Results),
plus a clean 4-stage fingerprint table for Strategy B.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT_A  = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_ACCENT_B  = RGBColor(0xFF, 0x6B, 0x6B)   # coral
C_ACCENT_C  = RGBColor(0x06, 0xD6, 0xA0)   # green
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
C_DIM       = RGBColor(0x88, 0x99, 0xAA)
C_DARK_CARD = RGBColor(0x1A, 0x2E, 0x42)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FIG_PIPELINE = "/tmp/figs/fig_pipeline_overview.png"
FIG_SPEEDUP  = "/tmp/figs/fig1_speedup_bar.png"

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    return s


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or C_WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)


def accent_bar(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), color)


def slide_number(slide, n, color):
    add_text(slide, str(n), Inches(12.7), Inches(7.15), Inches(0.5), Inches(0.3),
             size=10, color=color, align=PP_ALIGN.RIGHT)


def speaker_notes(slide, lines):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    first = True
    for ln in lines:
        p = notes.paragraphs[0] if first else notes.add_paragraph()
        first = False
        p.text = "• " + ln


def add_bullets(slide, items, x, y, w, h, size=24, color=None):
    """Big, breathable bullets. Items must already be ≤5 words each."""
    color = color or C_WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color


def assert_5x5(title, bullets):
    """Enforce 5x5 rule at build time."""
    assert len(bullets) <= 5, f"too many bullets ({len(bullets)}) on slide '{title}'"
    for b in bullets:
        n = len(b.split())
        assert n <= 6, f"bullet too long ({n} words) on '{title}': '{b}'"


# ── Slide factories ────────────────────────────────────────────────────────

def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, Inches(2.8), Inches(4.44), Inches(0.05), C_ACCENT_A)
    add_rect(slide, Inches(4.44), Inches(2.8), Inches(4.44), Inches(0.05), C_ACCENT_B)
    add_rect(slide, Inches(8.88), Inches(2.8), Inches(4.45), Inches(0.05), C_ACCENT_C)

    add_text(slide, "ML Experiment Optimization via Intelligent Caching",
             Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.2),
             size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Bachelor Thesis  |  Lennart Gorzel  |  IMC FH Krems",
             Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.5),
             size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "Three storytelling strategies  •  5×5 rule  •  7 minutes",
             Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.4),
             size=12, color=C_DIM, align=PP_ALIGN.CENTER, italic=True)

    cards = [
        ("Strategy A", "The Researcher's Day", C_ACCENT_A),
        ("Strategy B", "The Trust Problem",    C_ACCENT_B),
        ("Strategy C", "Numbers-First",        C_ACCENT_C),
    ]
    for i, (tag, sub, col) in enumerate(cards):
        bx = Inches(0.5 + i * 4.27)
        add_rect(slide, bx, Inches(3.8), Inches(3.9), Inches(1.5), C_DARK_CARD)
        add_text(slide, tag, bx + Inches(0.15), Inches(3.95), Inches(3.6), Inches(0.5),
                 size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, sub, bx + Inches(0.15), Inches(4.55), Inches(3.6), Inches(0.5),
                 size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    speaker_notes(slide, [
        "Pick ONE strategy for the actual defense; do not present all three.",
        "All three strategies are 7 minutes total and follow the same outline.",
        "Speaker notes for each slide live in this notes pane.",
    ])


def make_divider(prs, label, subtitle, color, n, timing):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, color)
    add_rect(slide, Inches(0.15), Inches(2.9), SLIDE_W - Inches(0.15), Inches(0.04), color)
    add_text(slide, label, Inches(1), Inches(1.4), Inches(11.3), Inches(1.2),
             size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, Inches(1), Inches(2.8), Inches(11.3), Inches(0.6),
             size=18, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)
    add_text(slide, timing, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
             size=11, color=C_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "7 min  •  CS committee  •  English  •  5×5 rule",
             Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
             size=12, color=C_DIM, align=PP_ALIGN.CENTER)
    slide_number(slide, n, color)


def make_slide(prs, section_tag, title, bullets, notes, color, n,
               figure=None, time_budget=None, fig_caption=None,
               fingerprint_table=False):
    """One content slide. Bullets enforced 5x5. Optional figure or table on the right."""
    assert_5x5(title, bullets)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    accent_bar(slide, color)

    add_text(slide, section_tag, Inches(0.35), Inches(0.15), Inches(8.5), Inches(0.3),
             size=11, bold=True, color=color)
    if time_budget:
        add_text(slide, time_budget, Inches(11.5), Inches(0.15), Inches(1.6), Inches(0.3),
                 size=11, color=color, align=PP_ALIGN.RIGHT, bold=True)

    add_text(slide, title, Inches(0.35), Inches(0.5), Inches(12.6), Inches(0.9),
             size=34, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.35), Inches(1.4), Inches(12.6), Inches(0.025), color)

    has_visual = figure is not None or fingerprint_table
    if has_visual:
        bullets_x, bullets_w = Inches(0.5),  Inches(5.5)
    else:
        bullets_x, bullets_w = Inches(1.5), Inches(10.3)

    add_bullets(slide, bullets, bullets_x, Inches(1.85), bullets_w, Inches(5.2),
                size=22 if has_visual else 26)

    if figure:
        fig_x, fig_y = Inches(6.5), Inches(1.7)
        fig_w, fig_h = Inches(6.6), Inches(5.0)
        add_rect(slide, fig_x, fig_y, fig_w, fig_h + Inches(0.1), C_WHITE)
        slide.shapes.add_picture(figure, fig_x + Inches(0.05), fig_y + Inches(0.05),
                                 width=fig_w - Inches(0.1))
        if fig_caption:
            add_text(slide, fig_caption, fig_x, Inches(6.9), fig_w, Inches(0.3),
                     size=10, color=C_DIM, align=PP_ALIGN.CENTER, italic=True)

    if fingerprint_table:
        draw_fingerprint_table(slide, color)

    speaker_notes(slide, notes)
    slide_number(slide, n, color)


def draw_fingerprint_table(slide, color):
    """Clean 4-stage fingerprint composition table on the right side of the slide."""
    x0, y0 = Inches(6.4), Inches(1.7)
    w_total, h_total = Inches(6.7), Inches(5.2)
    n_rows = 5
    n_cols = 2
    col_w = [Inches(2.3), Inches(4.4)]
    row_h = h_total / n_rows

    add_rect(slide, x0, y0, w_total, row_h, color)
    add_text(slide, "Stage", x0 + Inches(0.2), y0 + Inches(0.18), col_w[0], row_h,
             size=14, bold=True, color=C_WHITE)
    add_text(slide, "Fingerprint includes", x0 + col_w[0] + Inches(0.2),
             y0 + Inches(0.18), col_w[1], row_h,
             size=14, bold=True, color=C_WHITE)

    rows = [
        ("1. Preprocess",   "Raw EEG + filter config"),
        ("2. Features",     "Preprocessed + extraction config"),
        ("3. Selection",    "Features + threshold + top-K"),
        ("4. Model",        "Above + params + subject ID"),
    ]
    for i, (a, b) in enumerate(rows, start=1):
        y = y0 + row_h * i
        bg = C_DARK_CARD if i % 2 == 1 else C_BG
        add_rect(slide, x0, y, w_total, row_h, bg)
        is_last = i == len(rows)
        text_color = color if is_last else C_WHITE
        bold = is_last
        add_text(slide, a, x0 + Inches(0.2), y + Inches(0.22), col_w[0], row_h,
                 size=13, color=text_color, bold=bold)
        add_text(slide, b, x0 + col_w[0] + Inches(0.2), y + Inches(0.22),
                 col_w[1], row_h, size=13, color=text_color, bold=bold)

    add_text(slide, "Subject ID in Stage 4  →  fold identity guards the cache",
             x0, y0 + h_total + Inches(0.1), w_total, Inches(0.4),
             size=11, italic=True, color=color, align=PP_ALIGN.CENTER)


# ── Content ────────────────────────────────────────────────────────────────

STRATEGY_A = [
    dict(title="Nine Hours. Again.",
         bullets=[
             "128 subjects, 2,304 runs",
             "One run: 9–12 hours",
             "Change one parameter, restart",
             "80% redundant recomputation",
             "Bottleneck: pipeline, not model",
         ],
         notes=[
             "Open verbatim: 'I run an experiment. Nine hours. I change one number. Another nine hours.'",
             "Let the 2,304 number land. Pause.",
             "Don't define LOSO yet — say 'standard cross-validation'.",
             "Close: 'This isn't a model problem. It's a pipeline problem.'",
         ],
         time="1:30"),
    dict(title="The Parts Exist",
         bullets=[
             "SHA-256 hashing exists",
             "Caching libraries exist",
             "Model serialization exists",
             "None covers cross-validation",
             "Gap: assembling them safely",
         ],
         notes=[
             "Three lines, one per tool. Then the gap.",
             "Do NOT mention your combination. Just point at the gap.",
             "Transition: 'The parts exist. Assembling them safely doesn't.'",
         ],
         time="0:45"),
    dict(title="Identity of a Computation",
         bullets=[
             "All inputs → hash → key",
             "Same inputs, same key",
             "Different inputs, recompute",
             "Stages chained hierarchically",
             "Partial invalidation possible",
         ],
         notes=[
             "Anchor word: 'identity'. Return to it.",
             "Say 'fingerprint' instead of 'hash function' for the non-crypto folks.",
             "Transition: 'That's the concept. Now the wiring.'",
         ],
         time="1:00"),
    dict(title="Four-Stage Pipeline",
         bullets=[
             "Feature data feeds strategy",
             "LOSO split per subject",
             "Fingerprint decides cache hit",
             "Subject ID guards leakage",
             "Hierarchical, partial invalidation",
         ],
         notes=[
             "Walk through the figure left to right.",
             "Stop at 'Generate Fingerprint' — that's the novelty.",
             "Mention: change Stage 3 config → only 3+4 recompute.",
             "No code, no file tree. Stay on the diagram.",
         ],
         time="1:15",
         figure=FIG_PIPELINE,
         fig_caption="Pipeline with fingerprint-based cache lookup"),
    dict(title="Cold vs Warm Run",
         bullets=[
             "XGBoost: 54× median speedup",
             "Random Forest: 12× speedup",
             "128 LOSO folds tested",
             "100% cache hit rate",
             "Zero leakage detected",
         ],
         notes=[
             "Read the speedups off the chart. Pause after 54×.",
             "Lead with speed, then correctness.",
             "Reframe: 'Nine experiments a day, not one.'",
         ],
         time="1:30",
         figure=FIG_SPEEDUP,
         fig_caption="Cold vs warm execution, 128 LOSO folds"),
    dict(title="The Bottleneck Is Gone",
         bullets=[
             "Fingerprint every computation",
             "Cache the stable parts",
             "Recompute only what changed",
             "Nine hours to one",
             "Iteration is unlocked",
         ],
         notes=[
             "Closing line: 'Iteration cost was the bottleneck. It isn't anymore.'",
             "Mention generalization in one breath: any LOSO/k-fold pipeline.",
             "Stop. Pause. Wait for questions.",
         ],
         time="1:00"),
]

STRATEGY_B = [
    dict(title="Speed Easy. Safety Hard.",
         bullets=[
             "2,304 runs per experiment",
             "9 hours each, 80% redundant",
             "LOSO demands data isolation",
             "Wrong cache: silent leakage",
             "Real problem: safe caching",
         ],
         notes=[
             "Open with the hook: 'Cross-validation has a rule: don't touch the held-out subject. Caching has a habit: it remembers everything.'",
             "Frame as a TRUST problem, not a SPEED problem.",
             "Transition: 'Can we cache safely?'",
         ],
         time="1:00"),
    dict(title="Building Blocks, No Glue",
         bullets=[
             "SHA-256 hashing exists",
             "Caching libraries exist",
             "Model serialization exists",
             "None encodes fold identity",
             "Gap: cross-validation boundary",
         ],
         notes=[
             "Keep fast. The gap is the payoff.",
             "Last two bullets are the contribution shape. Don't say more yet.",
             "Transition: 'The boundary between them under cross-validation is missing.'",
         ],
         time="0:45"),
    dict(title="Identity Includes the Fold",
         bullets=[
             "Hash of all inputs",
             "Subject ID inside the hash",
             "Each fold: unique key",
             "Leakage prevented by construction",
             "Guarantee, not runtime check",
         ],
         notes=[
             "Emphasize 'by construction' — mathematically guaranteed.",
             "Slow down here. This is the insight.",
             "Transition: 'Now the wiring.'",
         ],
         time="1:15"),
    dict(title="Four Stages, Each Fingerprinted",
         bullets=[
             "Stage 1 hashes raw signals",
             "Stage 2 hashes features",
             "Stage 3 hashes selection config",
             "Stage 4 adds subject ID",
             "Partial invalidation cascades",
         ],
         notes=[
             "Read the table left column, then right.",
             "Stop at Stage 4. Say 'this is where safety lives.'",
             "The caption line summarizes the entire thesis contribution.",
         ],
         time="1:15",
         fingerprint_table=True),
    dict(title="Fast AND Correct",
         bullets=[
             "XGBoost: 54× median speedup",
             "Random Forest: 12× speedup",
             "100% cache hit rate",
             "Zero leakage cases confirmed",
             "Both claims hold together",
         ],
         notes=[
             "Present speed first, then correctness — that's the pair this strategy promises.",
             "'Zero leakage' is the payoff for the trust framing. Land it.",
             "Don't rush. This slide proves two things at once.",
         ],
         time="1:45",
         figure=FIG_SPEEDUP,
         fig_caption="Speed gains, with correctness preserved across all 128 folds"),
    dict(title="Safe Caching, Under Cross-Validation",
         bullets=[
             "Caching: well understood",
             "Safe caching under LOSO: not",
             "Fold ID in the fingerprint",
             "Zero correctness compromise",
             "Generalizes to k-fold pipelines",
         ],
         notes=[
             "Open with the contrast — that's the contribution.",
             "Closing line, slowly: 'Caching, safely, under cross-validation.'",
             "Stop. Don't summarize the summary.",
         ],
         time="1:00"),
]

STRATEGY_C = [
    dict(title="Nine Hours. Then One.",
         bullets=[
             "128 subjects, 2,304 runs",
             "Each run: 9–12 hours",
             "80% identical recomputation",
             "Question: cache without errors?",
             "Answer: one design choice",
         ],
         notes=[
             "Open with the hook: 'Nine hours. Then one. Same data, same model, same correctness.'",
             "Plant the question early — every following slide answers it.",
             "Numbers-first: result is already on the table.",
         ],
         time="1:15"),
    dict(title="Close, But Not Quite",
         bullets=[
             "Hashing exists",
             "Caching frameworks exist",
             "Model serialization exists",
             "No unified LOSO scheme",
             "Stale fold reads unprotected",
         ],
         notes=[
             "Fast slide, 45 seconds.",
             "Don't linger — the audience knows an answer is coming.",
         ],
         time="0:45"),
    dict(title="Fingerprint Every Computation",
         bullets=[
             "Hash of all inputs",
             "Same key, cache hit",
             "Different key, recompute",
             "Stages chained downstream",
             "Partial invalidation cascades",
         ],
         notes=[
             "Tight, conceptual. No implementation details yet.",
             "Transition: 'That's the idea. Here's the wiring.'",
         ],
         time="1:00"),
    dict(title="Four-Stage Hierarchical Cache",
         bullets=[
             "Pipeline split into four",
             "Each stage hashed separately",
             "Subject ID at training",
             "Each fold: distinct entry",
             "Threshold change: stages three, four",
         ],
         notes=[
             "Walk through the diagram briefly.",
             "Highlight Stage 4 = safety guarantee.",
             "Transition: 'Does it deliver on the promise?'",
         ],
         time="1:15",
         figure=FIG_PIPELINE,
         fig_caption="Hierarchical cache with fold-aware fingerprints"),
    dict(title="The Numbers Hold Up",
         bullets=[
             "XGBoost: 54× median speedup",
             "Random Forest: 12× speedup",
             "100% cache hit rate",
             "Zero leakage events",
             "Promise from slide one delivered",
         ],
         notes=[
             "Echo the hook: 'Remember slide one — nine hours, then one.'",
             "The last bullet closes the story loop.",
         ],
         time="1:30",
         figure=FIG_SPEEDUP,
         fig_caption="Median speedup across 128 LOSO folds, two model families"),
    dict(title="One Design Choice",
         bullets=[
             "Fingerprint every computation",
             "Include fold ID inside",
             "Four cascading stages",
             "12–54× measured speedup",
             "Generalizes to k-fold pipelines",
         ],
         notes=[
             "Echo opening line as closing line.",
             "Then stop. Pause. Wait for questions.",
         ],
         time="1:15"),
]


# ── Build ──────────────────────────────────────────────────────────────────

SECTION_NAMES = ["Motivation", "State of the Art", "Model",
                 "Implementation", "Results", "Conclusion"]


def emit_strategy(prs, name, color, slides_data, start_n):
    tag_prefix = f"Strategy {name[0]} — {name[2:]}"
    for i, s in enumerate(slides_data):
        section = SECTION_NAMES[i]
        make_slide(prs,
                   section_tag=f"{tag_prefix}  |  {section}",
                   title=s["title"],
                   bullets=s["bullets"],
                   notes=s["notes"],
                   color=color,
                   n=start_n + i,
                   time_budget=s["time"],
                   figure=s.get("figure"),
                   fig_caption=s.get("fig_caption"),
                   fingerprint_table=s.get("fingerprint_table", False))


def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)

    make_divider(prs, "Strategy A — The Researcher's Day",
                 "First-person pain → unified idea → relief",
                 C_ACCENT_A, 2,
                 "Motivation 1:30  •  SOTA 0:45  •  Model 1:00  •  Implementation 1:15  •  Results 1:30  •  Conclusion 1:00")
    emit_strategy(prs, "A — The Researcher's Day", C_ACCENT_A, STRATEGY_A, start_n=3)

    make_divider(prs, "Strategy B — The Trust Problem",
                 "Correctness framing → gap → fast AND safe",
                 C_ACCENT_B, 9,
                 "Motivation 1:00  •  SOTA 0:45  •  Model 1:15  •  Implementation 1:15  •  Results 1:45  •  Conclusion 1:00")
    emit_strategy(prs, "B — The Trust Problem", C_ACCENT_B, STRATEGY_B, start_n=10)

    make_divider(prs, "Strategy C — Numbers-First",
                 "Open with the number → backfill → close the loop",
                 C_ACCENT_C, 16,
                 "Motivation 1:15  •  SOTA 0:45  •  Model 1:00  •  Implementation 1:15  •  Results 1:30  •  Conclusion 1:15")
    emit_strategy(prs, "C — Numbers-First", C_ACCENT_C, STRATEGY_C, start_n=17)

    out = "/home/user/BachlorThesis/presentation_strategies_v2.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
